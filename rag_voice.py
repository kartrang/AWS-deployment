from typing import List, Dict, Optional, Tuple
import os
import tempfile
from datetime import datetime
import uuid
import asyncio
import base64
from io import BytesIO
from PIL import Image

import streamlit as st
from dotenv import load_dotenv
from qdrant_client import QdrantClient
from qdrant_client.http import models
from qdrant_client.http.models import Distance, VectorParams
# langchain moved/renamed text splitter modules across versions.
# Try the common import paths and provide a clear error if neither exists.
try:
    from langchain.text_splitter import RecursiveCharacterTextSplitter
except Exception:
    try:
        from langchain.text_splitters import RecursiveCharacterTextSplitter
    except Exception as e:
        try:
            # Some distributions install the splitter as a separate top-level package
            from langchain_text_splitters import RecursiveCharacterTextSplitter
        except Exception:
            raise ImportError(
                "RecursiveCharacterTextSplitter not found. Please install a compatible `langchain` or `langchain_text_splitters` package."
            ) from e
from langchain_community.document_loaders import PyPDFLoader
from fastembed import TextEmbedding
from openai import AsyncOpenAI, OpenAI
from openai.helpers import LocalAudioPlayer
import types

# Some third-party packages (older code) expect the `tf.contrib` namespace
# which was removed in TensorFlow 2.x. Create a lightweight shim before
# importing such packages so they can access `tf.contrib.distributions`.
try:
    import tensorflow as _tf
    if not hasattr(_tf, 'contrib'):
        _tf.contrib = types.SimpleNamespace()
        try:
            import tensorflow_probability as _tfp
            # Prefer to assign the actual distributions module. If for any
            # reason that's not sufficient for downstream code, copy known
            # attributes (like MultivariateNormalDiag) explicitly.
            try:
                _tf.contrib.distributions = _tfp.distributions
                # Ensure MultivariateNormalDiag exists on the target
                if not hasattr(_tf.contrib.distributions, 'MultivariateNormalDiag') and hasattr(_tfp.distributions, 'MultivariateNormalDiag'):
                    setattr(_tf.contrib.distributions, 'MultivariateNormalDiag', getattr(_tfp.distributions, 'MultivariateNormalDiag'))
            except Exception:
                # Fallback: create a namespace and copy attributes
                mod = types.SimpleNamespace()
                for name in dir(_tfp.distributions):
                    if not name.startswith('_'):
                        try:
                            setattr(mod, name, getattr(_tfp.distributions, name))
                        except Exception:
                            pass
                _tf.contrib.distributions = mod
        except Exception:
            # Provide an empty namespace as a fallback to avoid AttributeError
            _tf.contrib.distributions = types.SimpleNamespace()
except Exception:
    # If tensorflow isn't available or fails to import, continue and let
    # the downstream import raise a clearer error when needed.
    pass

# Note: The legacy 'agents' package (RL training lib) conflicts with code
# that expects OpenAI agents SDK. We'll stub Agent/Runner so the app can load
# with modes that don't use them (General Chat, Image Chat, Process Flow, etc.).
try:
    from agents import Agent, Runner
except Exception:
    # Fallback: create stub classes that allow the app to load even if
    # Agent/Runner aren't available. RAG mode will be limited but other
    # modes will work.
    class Agent:
        def __init__(self, *args, **kwargs):
            raise NotImplementedError("Agent class not available. This feature requires openai-agents SDK to be properly installed.")
    
    class Runner:
        @staticmethod
        def run(*args, **kwargs):
            raise NotImplementedError("Runner class not available. This feature requires openai-agents SDK to be properly installed.")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from docx import Document as DocxDocument
from elevenlabs import ElevenLabs, Voice, VoiceSettings
from pydub import AudioSegment

load_dotenv()

# Constants
COLLECTION_NAME = "voice-rag-agent"

def init_session_state() -> None:
    """Initialize Streamlit session state with default values."""
    defaults = {
        "initialized": False,
        "qdrant_url": "",
        "qdrant_api_key": "",
        "openai_api_key": "",
        "setup_complete": False,
        "client": None,
        "embedding_model": None,
        "processor_agent": None,
        "tts_agent": None,
        "selected_voice": "coral",
        "processed_documents": [],
        "mode": "General Chat",  # Default mode
        "selected_model": "gpt-4o",  # Default model
        "chat_history": [],
        "current_image": None,  # For Image Chat mode
        "image_chat_history": [],  # Image chat conversation history
        "image_base64": None,  # Base64 encoded image for API
        "process_flow_chat_history": [],  # Chat history for process flow mode
        "generated_flow_path": None,  # Path to generated process flow PPT
        "generated_excel_path": None,  # Path to generated process flow Excel
        "bpmn_chat_history": [],  # Chat history for BPMN mode
        "generated_bpmn_xml": None,  # Generated BPMN XML content
        "voice_file": None,  # Uploaded voice file for cloning
        "cloned_voice_id": None,  # ElevenLabs cloned voice ID
        "elevenlabs_api_key": "",  # ElevenLabs API key
        "voice_stability": 0.5,  # Voice stability (0-1)
        "voice_similarity": 0.75,  # Voice similarity boost (0-1)
        "voice_style": 0.0,  # Voice style exaggeration (0-1)
        "use_speaker_boost": True  # Enable speaker boost
    }
    
    for key, value in defaults.items():
        if key not in st.session_state:
            st.session_state[key] = value

def setup_sidebar() -> None:
    """Configure sidebar with API settings and voice options."""
    with st.sidebar:
        st.title("🔑 Configuration")
        st.markdown("---")
        
        # Mode selection
        st.markdown("### 🤖 Mode")
        mode_options = ["General Chat", "RAG (Document Q&A)", "Image Chat", "Process Flow Creator", "Voice Cloning", "BPMN Diagram Generator"]
        st.session_state.mode = st.radio(
            "Select Mode",
            options=mode_options,
            index=mode_options.index(st.session_state.mode) if st.session_state.mode in mode_options else 0,
            help="General Chat: Ask any question like ChatGPT\nRAG: Query uploaded documents\nImage Chat: Analyze and discuss images\nProcess Flow Creator: Generate process flow diagrams\nVoice Cloning: Clone your voice and convert text to speech\nBPMN Diagram Generator: Generate Business Process Model and Notation diagrams"
        )
        
        # Clear mode-specific session state when switching modes
        if "previous_mode" not in st.session_state:
            st.session_state.previous_mode = st.session_state.mode
        
        if st.session_state.previous_mode != st.session_state.mode:
            # Mode changed - clear mode-specific state
            if st.session_state.previous_mode == "BPMN Diagram Generator":
                # Clear BPMN-specific state
                st.session_state.bpmn_chat_history = []
                st.session_state.generated_bpmn_xml = None
            elif st.session_state.previous_mode == "Process Flow Creator":
                # Clear process flow state
                st.session_state.process_flow_chat_history = []
                st.session_state.generated_flow_path = None
                st.session_state.generated_excel_path = None
            elif st.session_state.previous_mode == "Image Chat":
                # Clear image chat state
                st.session_state.current_image = None
                st.session_state.image_chat_history = []
                st.session_state.image_base64 = None
            elif st.session_state.previous_mode == "Voice Cloning":
                # Clear voice cloning state
                st.session_state.voice_file = None
                st.session_state.cloned_voice_id = None
            
            # Update previous mode
            st.session_state.previous_mode = st.session_state.mode
        
        st.markdown("---")
        st.markdown("### 🧠 Model Selection")
        
        # Model categories
        openai_models = ["gpt-4o", "gpt-4o-mini", "gpt-4-turbo", "gpt-3.5-turbo"]
        opensource_models = ["llama-3.1-70b (via Groq)", "mixtral-8x7b (via Groq)", "gemma-7b (via Groq)"]
        
        all_models = openai_models + opensource_models
        
        current_index = all_models.index(st.session_state.selected_model) if st.session_state.selected_model in all_models else 0
        st.session_state.selected_model = st.selectbox(
            "Choose Model",
            options=all_models,
            index=current_index,
            help="Select AI model for responses"
        )
        
        st.markdown("---")
        
        # API Keys - Qdrant only required for RAG mode
        if st.session_state.mode == "RAG (Document Q&A)":
            st.session_state.qdrant_url = st.text_input(
                "Qdrant URL",
                value=st.session_state.qdrant_url,
                type="password",
                help="Get from Qdrant Cloud: https://cloud.qdrant.io/"
            )
            st.session_state.qdrant_api_key = st.text_input(
                "Qdrant API Key",
                value=st.session_state.qdrant_api_key,
                type="password",
                help="Your Qdrant Cloud API key"
            )
        
        # ElevenLabs API Key for Voice Cloning mode
        if st.session_state.mode == "Voice Cloning":
            st.session_state.elevenlabs_api_key = st.text_input(
                "ElevenLabs API Key",
                value=st.session_state.elevenlabs_api_key,
                type="password",
                help="Get from ElevenLabs: https://elevenlabs.io/"
            )
        else:
            st.session_state.openai_api_key = st.text_input(
                "OpenAI API Key",
                value=st.session_state.openai_api_key,
                type="password",
                help="Required for OpenAI models and TTS"
            )
        
        st.markdown("---")
        
        # Voice Settings - different for Voice Cloning vs other modes
        if st.session_state.mode == "Voice Cloning":
            st.markdown("### 🎛️ Voice Quality Controls")
            st.session_state.voice_stability = st.slider(
                "Stability",
                min_value=0.0,
                max_value=1.0,
                value=st.session_state.voice_stability,
                step=0.05,
                help="Higher = more consistent, Lower = more expressive"
            )
            st.session_state.voice_similarity = st.slider(
                "Similarity Boost",
                min_value=0.0,
                max_value=1.0,
                value=st.session_state.voice_similarity,
                step=0.05,
                help="How closely to match the cloned voice"
            )
            st.session_state.voice_style = st.slider(
                "Style Exaggeration",
                min_value=0.0,
                max_value=1.0,
                value=st.session_state.voice_style,
                step=0.05,
                help="Amplify the speaking style of the voice"
            )
            st.session_state.use_speaker_boost = st.checkbox(
                "Speaker Boost",
                value=st.session_state.use_speaker_boost,
                help="Enhance voice clarity and quality"
            )
        else:
            st.markdown("### 🎤 Voice Settings")
            voices = ["alloy", "ash", "ballad", "coral", "echo", "fable", "onyx", "nova", "sage", "shimmer", "verse"]
            st.session_state.selected_voice = st.selectbox(
                "Select Voice",
                options=voices,
                index=voices.index(st.session_state.selected_voice),
                help="Choose the voice for the audio response"
            )

def setup_qdrant() -> Tuple[QdrantClient, TextEmbedding]:
    """Initialize Qdrant client and embedding model."""
    if not all([st.session_state.qdrant_url, st.session_state.qdrant_api_key]):
        raise ValueError("Qdrant credentials not provided")
    
    try:
        client = QdrantClient(
            url=st.session_state.qdrant_url,
            api_key=st.session_state.qdrant_api_key,
            verify=True,  # Verify SSL certificates
            timeout=30,    # 30 second timeout
            check_compatibility=False  # Skip version compatibility check
        )
        
        # Test the connection by trying to get collections
        try:
            collections = client.get_collections()
            st.success(f"✅ Successfully connected to Qdrant! Found {len(collections.collections)} collections.")
        except Exception as conn_error:
            raise ValueError(f"Failed to connect to Qdrant: {str(conn_error)}. Please check your URL and API key.")
            
    except Exception as e:
        error_msg = str(e)
        if "404" in error_msg:
            raise ValueError("❌ Qdrant URL not found (404). Please check your Qdrant Cloud URL - it should look like: https://xyz-abc.us-east-1-0.aws.cloud.qdrant.io:6333")
        elif "401" in error_msg or "403" in error_msg:
            raise ValueError("❌ Qdrant authentication failed. Please check your API key.")
        else:
            raise ValueError(f"❌ Qdrant connection error: {error_msg}")
    
    embedding_model = TextEmbedding()
    test_embedding = list(embedding_model.embed(["test"]))[0]
    embedding_dim = len(test_embedding)
    
    try:
        client.create_collection(
            collection_name=COLLECTION_NAME,
            vectors_config=VectorParams(
                size=embedding_dim,
                distance=Distance.COSINE
            )
        )
    except Exception as e:
        if "already exists" not in str(e):
            raise e
    
    return client, embedding_model

def process_pdf(file) -> List:
    """Process PDF file and split into chunks with metadata."""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tmp_file:
            tmp_file.write(file.getvalue())
            loader = PyPDFLoader(tmp_file.name)
            documents = loader.load()
            
            # Add source metadata
            for doc in documents:
                doc.metadata.update({
                    "source_type": "pdf",
                    "file_name": file.name,
                    "timestamp": datetime.now().isoformat()
                })
            
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000,
                chunk_overlap=200
            )
            return text_splitter.split_documents(documents)
    except Exception as e:
        st.error(f"📄 PDF processing error: {str(e)}")
        return []

def extract_text_from_docx(file) -> str:
    """Extract text from Word document."""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as tmp_file:
            tmp_file.write(file.getvalue())
            doc = DocxDocument(tmp_file.name)
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            return "\n".join(full_text)
    except Exception as e:
        st.error(f"📄 Word document processing error: {str(e)}")
        return ""

def extract_text_from_pptx(file) -> str:
    """Extract text from PowerPoint presentation."""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp_file:
            tmp_file.write(file.getvalue())
            prs = Presentation(tmp_file.name)
            full_text = []
            for slide_num, slide in enumerate(prs.slides, 1):
                full_text.append(f"\n--- Slide {slide_num} ---")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        full_text.append(shape.text)
            return "\n".join(full_text)
    except Exception as e:
        st.error(f"📊 PowerPoint processing error: {str(e)}")
        return ""

def process_presentation_docs(files) -> str:
    """Process multiple documents (PDF, Word, PPT) and extract text."""
    all_content = []
    for file in files:
        file_ext = file.name.lower().split('.')[-1]
        st.info(f"Processing {file.name}...")
        
        if file_ext == 'pdf':
            docs = process_pdf(file)
            content = "\n".join([doc.page_content for doc in docs])
            all_content.append(f"\n=== {file.name} ===\n{content}")
        elif file_ext == 'docx':
            content = extract_text_from_docx(file)
            all_content.append(f"\n=== {file.name} ===\n{content}")
        elif file_ext in ['ppt', 'pptx']:
            content = extract_text_from_pptx(file)
            all_content.append(f"\n=== {file.name} ===\n{content}")
        else:
            st.warning(f"Unsupported file type: {file_ext}")
    
    return "\n\n".join(all_content)

def store_embeddings(
    client: QdrantClient,
    embedding_model: TextEmbedding,
    documents: List,
    collection_name: str
) -> None:
    """Store document embeddings in Qdrant."""
    for doc in documents:
        embedding = list(embedding_model.embed([doc.page_content]))[0]
        # Convert numpy array to list if needed
        vector = embedding.tolist() if hasattr(embedding, 'tolist') else list(embedding)
        client.upsert(
            collection_name=collection_name,
            points=[
                models.PointStruct(
                    id=str(uuid.uuid4()),
                    vector=vector,
                    payload={
                        "content": doc.page_content,
                        **doc.metadata
                    }
                )
            ]
        )

def setup_agents(openai_api_key: str) -> Tuple[Agent, Agent]:
    """Initialize the processor and TTS agents."""
    os.environ["OPENAI_API_KEY"] = openai_api_key
    
    processor_agent = Agent(
        name="Documentation Processor",
        instructions="""You are a helpful documentation assistant. Your task is to:
        1. Analyze the provided documentation content
        2. Answer the user's question clearly and concisely
        3. Include relevant examples when available
        4. Cite the source files when referencing specific content
        5. Keep responses natural and conversational
        6. Format your response in a way that's easy to speak out loud""",
        model="gpt-4o"
    )

    tts_agent = Agent(
        name="Text-to-Speech Agent",
        instructions="""You are a text-to-speech agent. Your task is to:
        1. Convert the processed documentation response into natural speech
        2. Maintain proper pacing and emphasis
        3. Handle technical terms clearly
        4. Keep the tone professional but friendly
        5. Use appropriate pauses for better comprehension
        6. Ensure the speech is clear and well-articulated""",
        model="gpt-4o"
    )
    
    return processor_agent, tts_agent

async def process_query(
    query: str,
    client: QdrantClient,
    embedding_model: TextEmbedding,
    collection_name: str,
    openai_api_key: str,
    voice: str
) -> Dict:
    """Process user query and generate voice response."""
    try:
        st.info("🔄 Step 1: Generating query embedding and searching documents...")
        # Get query embedding and search
        query_embedding = list(embedding_model.embed([query]))[0]
        st.write(f"Generated embedding of size: {len(query_embedding)}")
        
        # Convert numpy array to list if needed
        query_vector = query_embedding.tolist() if hasattr(query_embedding, 'tolist') else list(query_embedding)
        
        search_response = client.query_points(
            collection_name=collection_name,
            query=query_vector,
            limit=3,
            with_payload=True
        )
        
        search_results = search_response.points if hasattr(search_response, 'points') else []
        st.write(f"Found {len(search_results)} relevant documents")
        
        if not search_results:
            raise Exception("No relevant documents found in the vector database")
        
        st.info("🔄 Step 2: Preparing context from search results...")
        # Prepare context from search results
        context = "Based on the following documentation:\n\n"
        for i, result in enumerate(search_results, 1):
            payload = result.payload
            if not payload:
                continue
            content = payload.get('content', '')
            source = payload.get('file_name', 'Unknown Source')
            context += f"From {source}:\n{content}\n\n"
            st.write(f"Document {i} from: {source}")
        
        context += f"\nUser Question: {query}\n\n"
        context += "Please provide a clear, concise answer that can be easily spoken out loud."
        
        st.info("🔄 Step 3: Generating response with OpenAI...")
        # Use OpenAI API directly to generate response from context
        client = OpenAI(api_key=openai_api_key)
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[
                {
                    "role": "system",
                    "content": "You are a helpful documentation assistant. Answer the user's question based on the provided documents. Be clear, concise, and easy to understand. Format your response as natural speech."
                },
                {
                    "role": "user",
                    "content": context
                }
            ],
            max_tokens=500
        )
        text_response = response.choices[0].message.content
        if not text_response:
            raise Exception("No response generated from the model")
        st.write(f"Generated text response of length: {len(text_response)}")
        
        st.info("🔄 Step 4: Generating audio...")
        # Generate audio response
        async_openai = AsyncOpenAI(api_key=openai_api_key)
        
        try:
            # Try to use LocalAudioPlayer for real-time playback if available
            try:
                from openai.helpers import LocalAudioPlayer
                
                async with async_openai.audio.speech.with_streaming_response.create(
                    model="tts-1",
                    voice=voice,
                    input=text_response,
                    response_format="pcm",
                ) as stream_response:
                    st.write("🎵 Playing audio (real-time)...")
                    await LocalAudioPlayer().play(stream_response)
                    st.success("✅ Audio playback complete")
            except (ImportError, OSError) as e:
                st.warning(f"⚠️ Real-time audio playback not available in this environment: {str(e)}")
                st.info("💡 Generating downloadable MP3 instead...")
            
            # Always generate MP3 for download
            st.write("Generating downloadable MP3 version...")
            audio_response = await async_openai.audio.speech.create(
                model="tts-1",
                voice=voice,
                input=text_response,
                response_format="mp3"
            )
            
            temp_dir = tempfile.gettempdir()
            audio_path = os.path.join(temp_dir, f"response_{uuid.uuid4()}.mp3")
            
            with open(audio_path, "wb") as f:
                f.write(audio_response.content)
            st.success(f"✅ MP3 file generated successfully")
            
        except Exception as audio_error:
            st.error(f"❌ Audio generation failed: {str(audio_error)}")
            return {
                "status": "error",
                "error": f"Audio generation failed: {str(audio_error)}",
                "query": query
            }
        
        st.success("✅ Query processing complete!")
        return {
            "status": "success",
            "text_response": text_response,
            "audio_path": audio_path,
            "sources": [r.payload.get('file_name', 'Unknown Source') for r in search_results if r.payload]
        }
    
    except Exception as e:
        st.error(f"❌ Error during query processing: {str(e)}")
        return {
            "status": "error",
            "error": str(e),
            "query": query
        }

async def generate_process_flow(
    process_description: str,
    openai_api_key: str
) -> str:
    """Generate process flow diagram in PowerPoint from process description."""
    try:
        st.info("🔄 Step 1: Analyzing process flow with AI...")
        
        # Use OpenAI to extract structured process flow
        client = OpenAI(api_key=openai_api_key)
        
        prompt = f"""Analyze the following process description and create a structured SWIMLANE process flow diagram.

PROCESS DESCRIPTION:
{process_description}

Extract:
1. Process title/name
2. All roles/departments involved (swimlanes)
3. All steps in sequence with role assignments
4. Flow connections between steps (including cross-lane handoffs)

Format your response as JSON with this structure:
{{
    "title": "Process Name",
    "lanes": [
        {{
            "id": "lane1",
            "name": "Role/Department Name"
        }}
    ],
    "steps": [
        {{
            "id": 1,
            "lane_id": "lane1",
            "type": "start|process|decision|end",
            "text": "Step description",
            "tool": "Tool/System name (if applicable)",
            "next": [2]  // IDs of next steps
        }}
    ]
}}

Step types:
- "start": Beginning of process (oval)
- "process": Regular activity/task (rectangle)
- "decision": Decision point with Yes/No branches (diamond)
- "end": End of process (oval)

Important:
- Identify ALL roles/departments involved and create lanes for each
- Assign EVERY step to the appropriate lane_id
- Include tools/systems used in each step
- Show handoffs between roles/departments through connections
- For decision points, specify both branches in the "next" array

Example roles: Employee, Manager, HR Department, System, Customer, etc.
"""
        
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"}
        )
        
        import json
        flow_data = json.loads(response.choices[0].message.content or "{}")
        st.success(f"✅ Process flow analyzed! Found {len(flow_data.get('steps', []))} steps")
        
        st.info("🔄 Step 2: Creating process flow diagram...")
        
        # Create presentation
        from pptx.util import Inches, Pt
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.dml.color import RGBColor
        
        prs = Presentation()
        prs.slide_width = Inches(13.33)  # Widescreen
        prs.slide_height = Inches(7.5)
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        if title:
            title.text = flow_data.get("title", "Process Flow Diagram")
        
        # Add blank slide for flowchart
        blank_layout = prs.slide_layouts[6]  # Blank layout
        flow_slide = prs.slides.add_slide(blank_layout)
        
        # Add title to flow slide
        title_box = flow_slide.shapes.add_textbox(
            Inches(0.5), Inches(0.3), Inches(12), Inches(0.6)
        )
        title_frame = title_box.text_frame
        title_frame.text = flow_data.get("title", "Process Flow")
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(28)
        title_para.font.bold = True
        title_para.alignment = PP_ALIGN.CENTER
        
        # Draw swimlane diagram
        lanes = flow_data.get("lanes", [])
        steps = flow_data.get("steps", [])
        
        # If no lanes defined, create a fallback
        if not lanes:
            st.warning("⚠️ No roles/lanes detected. Creating single-lane diagram.")
            lanes = [{"id": "default", "name": "Process Flow"}]
            for step in steps:
                step["lane_id"] = "default"
        
        step_shapes = {}
        lane_positions = {}
        
        # Swimlane layout parameters
        left_margin = Inches(0.5)
        top_margin = Inches(1.2)
        label_width = Inches(1.5)
        lane_height = Inches(1.4)
        diagram_width = Inches(11.5)
        step_width = Inches(1.8)
        step_height = Inches(0.7)
        horizontal_spacing = Inches(2.2)
        
        # Lane colors
        lane_colors = [
            RGBColor(68, 114, 196),   # Blue
            RGBColor(112, 173, 71),    # Green
            RGBColor(255, 192, 0),     # Orange
            RGBColor(237, 125, 49),    # Red-orange
            RGBColor(158, 72, 162),    # Purple
            RGBColor(37, 170, 226),    # Cyan
        ]
        
        # Draw swimlanes
        for idx, lane in enumerate(lanes):
            lane_id = lane.get("id")
            lane_name = lane.get("name", f"Lane {idx+1}")
            color = lane_colors[idx % len(lane_colors)]
            
            y_pos = top_margin + (idx * lane_height)
            lane_positions[lane_id] = {
                "y_start": y_pos,
                "y_end": y_pos + lane_height,
                "y_center": y_pos + (lane_height / 2)
            }
            
            # Draw lane header (label)
            header = flow_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left_margin,
                y_pos,
                label_width,
                lane_height
            )
            header.fill.solid()
            header.fill.fore_color.rgb = color
            header.line.color.rgb = RGBColor(50, 50, 50)
            header.line.width = Pt(1)
            
            # Add lane label text
            text_frame = header.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = lane_name
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            text_frame.vertical_anchor = 1  # Middle
            text_frame.word_wrap = True
            
            # Draw lane body (background)
            body = flow_slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left_margin + label_width,
                y_pos,
                diagram_width,
                lane_height
            )
            body.fill.background()
            body.line.color.rgb = RGBColor(150, 150, 150)
            body.line.width = Pt(1)
        
        # Group steps by lane and calculate positions
        lane_step_counts = {}
        for step in steps:
            lane_id = step.get("lane_id", "default")
            lane_step_counts[lane_id] = lane_step_counts.get(lane_id, 0) + 1
        
        lane_step_indices = {lane_id: 0 for lane_id in lane_step_counts.keys()}
        
        # Position and draw process steps
        for step in steps:
            step_id = step.get("id")
            lane_id = step.get("lane_id", "default")
            step_type = step.get("type", "process")
            step_text = step.get("text", "")
            tool = step.get("tool", "")
            
            if lane_id not in lane_positions:
                continue
            
            # Calculate x position based on step sequence
            step_index = lane_step_indices[lane_id]
            lane_step_indices[lane_id] += 1
            
            x_pos = left_margin + label_width + Inches(0.5) + (step_index * horizontal_spacing)
            y_center = lane_positions[lane_id]["y_center"]
            y_pos = y_center - (step_height / 2)
            
            # Create shape based on type
            if step_type == "start" or step_type == "end":
                shape = flow_slide.shapes.add_shape(
                    MSO_SHAPE.OVAL,
                    x_pos, y_pos, step_width, step_height
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(146, 208, 80)  # Green
            elif step_type == "decision":
                shape = flow_slide.shapes.add_shape(
                    MSO_SHAPE.DIAMOND,
                    x_pos, y_pos, step_width, step_height
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(255, 217, 102)  # Yellow
            else:
                shape = flow_slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    x_pos, y_pos, step_width, step_height
                )
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(189, 215, 238)  # Light blue
            
            shape.line.color.rgb = RGBColor(0, 0, 0)
            shape.line.width = Pt(1.5)
            
            # Add step text
            text_frame = shape.text_frame
            text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = step_text
            p.font.size = Pt(9)
            p.alignment = PP_ALIGN.CENTER
            text_frame.word_wrap = True
            text_frame.vertical_anchor = 1  # Middle
            
            # Add tool info if available
            if tool:
                p2 = text_frame.add_paragraph()
                p2.text = f"[{tool}]"
                p2.font.size = Pt(7)
                p2.font.italic = True
                p2.alignment = PP_ALIGN.CENTER
            
            # Store shape reference
            step_shapes[step_id] = shape
        
        # Draw connectors between steps
        from pptx.enum.shapes import MSO_CONNECTOR
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        
        for step in steps:
            step_id = step.get("id")
            next_steps = step.get("next", [])
            
            if step_id in step_shapes:
                current_shape = step_shapes[step_id]
                
                for next_id in next_steps:
                    if next_id in step_shapes:
                        next_shape = step_shapes[next_id]
                        
                        # Calculate connector start and end points
                        start_x = current_shape.left + current_shape.width
                        start_y = current_shape.top + (current_shape.height // 2)
                        end_x = next_shape.left
                        end_y = next_shape.top + (next_shape.height // 2)
                        
                        # Draw connector
                        connector = flow_slide.shapes.add_connector(
                            MSO_CONNECTOR.STRAIGHT,
                            start_x, start_y,
                            end_x, end_y
                        )
                        connector.line.color.rgb = RGBColor(0, 0, 0)
                        connector.line.width = Pt(1.5)
                        connector.line.dash_style = MSO_LINE_DASH_STYLE.SOLID
        
        # Save presentation
        temp_dir = tempfile.gettempdir()
        output_path = os.path.join(temp_dir, f"process_flow_{uuid.uuid4()}.pptx")
        prs.save(output_path)
        
        st.success(f"✅ Process flow diagram created with {len(steps)} steps!")
        return output_path
        
    except Exception as e:
        st.error(f"❌ Error generating process flow: {str(e)}")
        raise e

async def generate_process_flow_excel(flow_data: dict) -> str:
    """Generate process flow diagram in Excel from flow data."""
    try:
        import openpyxl
        from openpyxl.styles import PatternFill, Border, Side, Alignment, Font
        from openpyxl.utils import get_column_letter
        
        st.info("🔄 Creating Excel swimlane diagram...")
        
        # Create workbook
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Process Flow"
        
        # Remove gridlines
        ws.sheet_view.showGridLines = False
        
        # Get data
        lanes = flow_data.get("lanes", [])
        steps = flow_data.get("steps", [])
        
        # If no lanes, create default
        if not lanes:
            lanes = [{"id": "default", "name": "Process Flow"}]
            for step in steps:
                step["lane_id"] = "default"
        
        # Define colors (RGB format)
        lane_colors = [
            "4472C4",  # Blue
            "70AD47",  # Green
            "FFC000",  # Orange
            "ED7D31",  # Red-orange
            "9E48A2",  # Purple
            "25AAE2",  # Cyan
        ]
        
        # Configuration
        start_row = 2
        start_col = 2
        label_width = 3  # columns for label
        lane_height = 4  # rows per lane
        step_width = 3   # columns per step
        step_height = 2  # rows per step
        
        # Set column widths
        for col in range(1, 50):
            ws.column_dimensions[get_column_letter(col)].width = 12
        
        # Set label column width
        for col in range(start_col, start_col + label_width):
            ws.column_dimensions[get_column_letter(col)].width = 18
        
        # Draw title
        ws.merge_cells(start_row=1, start_column=start_col, end_row=1, end_column=start_col + 20)
        title_cell = ws.cell(row=1, column=start_col)
        title_cell.value = flow_data.get("title", "Process Flow Diagram")
        title_cell.font = Font(size=18, bold=True, color="000000")
        title_cell.alignment = Alignment(horizontal='center', vertical='center')
        
        # Lane positions
        lane_row_positions = {}
        current_row = start_row
        
        # Border styles
        thin_border = Border(
            left=Side(style='thin'),
            right=Side(style='thin'),
            top=Side(style='thin'),
            bottom=Side(style='thin')
        )
        
        # Draw swimlanes
        for idx, lane in enumerate(lanes):
            lane_id = lane.get("id")
            lane_name = lane.get("name", f"Lane {idx+1}")
            color = lane_colors[idx % len(lane_colors)]
            
            lane_row_positions[lane_id] = {
                "start": current_row,
                "center": current_row + (lane_height // 2)
            }
            
            # Draw lane header (label)
            ws.merge_cells(
                start_row=current_row,
                start_column=start_col,
                end_row=current_row + lane_height - 1,
                end_column=start_col + label_width - 1
            )
            header_cell = ws.cell(row=current_row, column=start_col)
            header_cell.value = lane_name
            header_cell.fill = PatternFill(start_color=color, end_color=color, fill_type="solid")
            header_cell.font = Font(bold=True, color="FFFFFF", size=12)
            header_cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
            header_cell.border = thin_border
            
            # Draw lane body (background area)
            for r in range(current_row, current_row + lane_height):
                for c in range(start_col + label_width, start_col + 40):
                    cell = ws.cell(row=r, column=c)
                    cell.fill = PatternFill(start_color="F0F0F0", end_color="F0F0F0", fill_type="solid")
                    cell.border = Border(
                        left=Side(style='thin', color='C0C0C0'),
                        right=Side(style='thin', color='C0C0C0'),
                        top=Side(style='thin', color='C0C0C0'),
                        bottom=Side(style='thin', color='C0C0C0')
                    )
            
            current_row += lane_height
        
        # Draw process steps using global sequence for column positioning
        step_cell_positions = {}
        current_col = start_col + label_width + 1
        
        for step in steps:
            step_id = step.get("id")
            lane_id = step.get("lane_id", "default")
            step_type = step.get("type", "process")
            step_text = step.get("text", "")
            tool = step.get("tool", "")
            
            if lane_id not in lane_row_positions:
                continue
            
            # Determine color based on type
            if step_type == "start" or step_type == "end":
                fill_color = "92D050"  # Green
            elif step_type == "decision":
                fill_color = "FFD966"  # Yellow
            else:
                fill_color = "BDD7EE"  # Light blue
            
            # Calculate step position using global column sequence and lane-specific row
            lane_center_row = lane_row_positions[lane_id]["center"]
            step_start_row = lane_center_row - 1
            step_start_col = current_col
            
            # Store position for connectors
            step_cell_positions[step_id] = {
                "row": step_start_row,
                "col": current_col + (step_width // 2)
            }
            
            # Draw step box
            ws.merge_cells(
                start_row=step_start_row,
                start_column=step_start_col,
                end_row=step_start_row + step_height - 1,
                end_column=step_start_col + step_width - 1
            )
            step_cell = ws.cell(row=step_start_row, column=step_start_col)
            
            # Add text with tool info
            if tool:
                step_cell.value = f"{step_text}\n[{tool}]"
            else:
                step_cell.value = step_text
            
            step_cell.fill = PatternFill(start_color=fill_color, end_color=fill_color, fill_type="solid")
            step_cell.font = Font(size=10, bold=True if step_type in ["start", "end"] else False)
            step_cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=True)
            step_cell.border = Border(
                left=Side(style='medium'),
                right=Side(style='medium'),
                top=Side(style='medium'),
                bottom=Side(style='medium')
            )
            
            # Add arrow to next step (simple representation with →)
            next_steps = step.get("next", [])
            if next_steps and len(next_steps) > 0:
                arrow_col = step_start_col + step_width
                arrow_cell = ws.cell(row=step_start_row, column=arrow_col)
                arrow_cell.value = "→"
                arrow_cell.font = Font(size=14, bold=True)
                arrow_cell.alignment = Alignment(horizontal='center', vertical='center')
            
            # Advance column for next step (global sequence)
            current_col += step_width + 1
        
        # Set row heights for better appearance
        for row in range(1, current_row + 5):
            ws.row_dimensions[row].height = 25
        
        # Save Excel file
        temp_dir = tempfile.gettempdir()
        output_path = os.path.join(temp_dir, f"process_flow_{uuid.uuid4()}.xlsx")
        wb.save(output_path)
        
        st.success(f"✅ Excel diagram created with {len(steps)} steps!")
        return output_path
        
    except Exception as e:
        st.error(f"❌ Error generating Excel: {str(e)}")
        raise e

async def generate_bpmn_diagram(
    process_description: str,
    openai_api_key: str
) -> str:
    """Generate BPMN XML diagram from process description using OpenAI."""
    try:
        st.info("🔄 Step 1: Analyzing process with AI...")
        
        client = OpenAI(api_key=openai_api_key)
        
        prompt = f"""Analyze the following business process description and create a structured BPMN (Business Process Model and Notation) diagram.

PROCESS DESCRIPTION:
{process_description}

Extract and structure the process into BPMN elements. Return JSON with this structure:
{{
    "process_name": "Name of the process",
    "process_id": "ProcessID",
    "elements": [
        {{
            "id": "unique_id",
            "type": "startEvent|task|userTask|serviceTask|exclusiveGateway|parallelGateway|inclusiveGateway|endEvent",
            "name": "Element name/description",
            "incoming": ["flow_id_1"],
            "outgoing": ["flow_id_2"]
        }}
    ],
    "flows": [
        {{
            "id": "flow_id",
            "sourceRef": "source_element_id",
            "targetRef": "target_element_id",
            "name": "Condition (for gateway branches)"
        }}
    ]
}}

BPMN Element Types:
- startEvent: Single start point (no incoming flows)
- task: Generic task/activity
- userTask: Task performed by a human
- serviceTask: Automated/system task  
- exclusiveGateway: Decision point (one path chosen)
- parallelGateway: Parallel execution (all paths taken simultaneously)
- inclusiveGateway: Conditional parallel (one or more paths taken)
- endEvent: End point (no outgoing flows)

Rules:
1. Every process must have exactly ONE startEvent
2. Every process must have at least ONE endEvent
3. Use clear, descriptive names for all elements
4. For gateways, specify conditions in flow names
5. Ensure all IDs are unique
6. All elements must be connected via flows (except start has no incoming, end has no outgoing)

Example:
For "Customer orders product online":
- Start Event
- User Task: "Browse products"
- User Task: "Add to cart"
- User Task: "Checkout"
- Service Task: "Process payment"
- Exclusive Gateway: "Payment successful?"
  - Yes → Service Task: "Send confirmation"
  - No → End Event: "Order failed"
- Service Task: "Ship product"
- End Event: "Order complete"
"""
        
        response = client.chat.completions.create(
            model="gpt-4o",
            messages=[{"role": "user", "content": prompt}],
            response_format={"type": "json_object"}
        )
        
        import json
        bpmn_data = json.loads(response.choices[0].message.content or "{}")
        elements = bpmn_data.get("elements", [])
        flows = bpmn_data.get("flows", [])
        
        st.success(f"✅ Process analyzed! Found {len(elements)} elements and {len(flows)} flows")
        
        # Validate BPMN structure
        st.info("🔄 Step 2: Validating BPMN structure...")
        
        # Check for start events
        start_events = [e for e in elements if e.get("type") == "startEvent"]
        if len(start_events) == 0:
            st.error("❌ Validation Error: Process must have at least one Start Event")
            raise ValueError("Missing start event in process")
        elif len(start_events) > 1:
            st.warning(f"⚠️ Warning: Process has {len(start_events)} start events. BPMN best practice is to have exactly one.")
        
        # Check for end events
        end_events = [e for e in elements if e.get("type") == "endEvent"]
        if len(end_events) == 0:
            st.error("❌ Validation Error: Process must have at least one End Event")
            raise ValueError("Missing end event in process")
        
        # Check for unique IDs
        element_ids = [e.get("id") for e in elements]
        if len(element_ids) != len(set(element_ids)):
            st.error("❌ Validation Error: All element IDs must be unique")
            raise ValueError("Duplicate element IDs detected")
        
        # Check that flows reference valid elements
        element_id_set = set(element_ids)
        for flow in flows:
            source = flow.get("sourceRef")
            target = flow.get("targetRef")
            if source not in element_id_set:
                st.error(f"❌ Validation Error: Flow references unknown source '{source}'")
                raise ValueError(f"Invalid flow source: {source}")
            if target not in element_id_set:
                st.error(f"❌ Validation Error: Flow references unknown target '{target}'")
                raise ValueError(f"Invalid flow target: {target}")
        
        st.success("✅ BPMN structure validated!")
        
        st.info("🔄 Step 3: Generating BPMN XML...")
        
        # Generate BPMN XML
        import xml.etree.ElementTree as ET
        
        # Create root definitions element
        root = ET.Element("definitions")
        root.set("xmlns", "http://www.omg.org/spec/BPMN/20100524/MODEL")
        root.set("xmlns:bpmndi", "http://www.omg.org/spec/BPMN/20100524/DI")
        root.set("xmlns:dc", "http://www.omg.org/spec/DD/20100524/DC")
        root.set("xmlns:di", "http://www.omg.org/spec/DD/20100524/DI")
        root.set("xmlns:xsi", "http://www.w3.org/2001/XMLSchema-instance")
        root.set("id", "Definitions_1")
        root.set("targetNamespace", "http://bpmn.io/schema/bpmn")
        
        # Create process element
        process = ET.SubElement(root, "process")
        process_id = bpmn_data.get("process_id", "Process_1")
        process.set("id", process_id)
        process.set("name", bpmn_data.get("process_name", "Business Process"))
        process.set("isExecutable", "false")
        
        # Add all elements
        for elem in elements:
            elem_type = elem.get("type", "task")
            # Generate ID if not present and store it back in the elem dict
            if "id" not in elem or not elem["id"]:
                elem["id"] = f"Activity_{uuid.uuid4().hex[:8]}"
            elem_id = elem["id"]
            elem_name = elem.get("name", "")
            
            # Create element
            bpmn_elem = ET.SubElement(process, elem_type)
            bpmn_elem.set("id", elem_id)
            if elem_name:
                bpmn_elem.set("name", elem_name)
            
            # Add incoming flows
            for incoming in elem.get("incoming", []):
                inc = ET.SubElement(bpmn_elem, "incoming")
                inc.text = incoming
            
            # Add outgoing flows  
            for outgoing in elem.get("outgoing", []):
                out = ET.SubElement(bpmn_elem, "outgoing")
                out.text = outgoing
        
        # Add all flows (sequence flows)
        for flow in flows:
            seq_flow = ET.SubElement(process, "sequenceFlow")
            # Generate ID if not present and store it back in the flow dict
            if "id" not in flow or not flow["id"]:
                flow["id"] = f"Flow_{uuid.uuid4().hex[:8]}"
            flow_id = flow["id"]
            seq_flow.set("id", flow_id)
            seq_flow.set("sourceRef", flow.get("sourceRef", ""))
            seq_flow.set("targetRef", flow.get("targetRef", ""))
            if flow.get("name"):
                seq_flow.set("name", flow["name"])
        
        # Add BPMN Diagram Interchange (visual layout)
        st.info("🔄 Step 4: Adding diagram layout...")
        
        bpmn_diagram = ET.SubElement(root, "bpmndi:BPMNDiagram")
        bpmn_diagram.set("id", "BPMNDiagram_1")
        
        bpmn_plane = ET.SubElement(bpmn_diagram, "bpmndi:BPMNPlane")
        bpmn_plane.set("id", "BPMNPlane_1")
        bpmn_plane.set("bpmnElement", process_id)
        
        # Layout elements horizontally with spacing
        x_position = 100
        y_position = 100
        x_spacing = 200
        
        # Add shapes for each element
        for elem in elements:
            # Use the ID that was already assigned (no regeneration)
            elem_id = elem["id"]
            elem_type = elem.get("type", "task")
            
            shape = ET.SubElement(bpmn_plane, "bpmndi:BPMNShape")
            shape.set("id", f"{elem_id}_di")
            shape.set("bpmnElement", elem_id)
            
            bounds = ET.SubElement(shape, "dc:Bounds")
            bounds.set("x", str(x_position))
            bounds.set("y", str(y_position))
            
            # Set dimensions based on element type
            if "Event" in elem_type:
                # Events are circular (36x36)
                bounds.set("width", "36")
                bounds.set("height", "36")
            elif "Gateway" in elem_type:
                # Gateways are diamonds (50x50)
                bounds.set("width", "50")
                bounds.set("height", "50")
            else:
                # Tasks are rectangles (100x80)
                bounds.set("width", "100")
                bounds.set("height", "80")
            
            # Move to next position
            x_position += x_spacing
        
        # Add edges for each sequence flow
        for flow in flows:
            # Use the ID that was already assigned (no regeneration)
            flow_id = flow["id"]
            source_ref = flow.get("sourceRef", "")
            target_ref = flow.get("targetRef", "")
            
            # Find source and target positions
            source_idx = next((i for i, e in enumerate(elements) if e.get("id") == source_ref), 0)
            target_idx = next((i for i, e in enumerate(elements) if e.get("id") == target_ref), 1)
            
            source_x = 100 + (source_idx * x_spacing) + 50
            source_y = y_position + 40
            target_x = 100 + (target_idx * x_spacing)
            target_y = y_position + 40
            
            edge = ET.SubElement(bpmn_plane, "bpmndi:BPMNEdge")
            edge.set("id", f"{flow_id}_di")
            edge.set("bpmnElement", flow_id)
            
            # Add waypoints (connection points)
            waypoint1 = ET.SubElement(edge, "di:waypoint")
            waypoint1.set("x", str(source_x))
            waypoint1.set("y", str(source_y))
            
            waypoint2 = ET.SubElement(edge, "di:waypoint")
            waypoint2.set("x", str(target_x))
            waypoint2.set("y", str(target_y))
        
        # Convert to string with proper formatting
        xml_string = ET.tostring(root, encoding="unicode", method="xml")
        
        # Pretty print the XML
        from xml.dom import minidom
        dom = minidom.parseString(xml_string)
        pretty_xml = dom.toprettyxml(indent="  ")
        
        # Remove extra blank lines
        pretty_xml = "\n".join([line for line in pretty_xml.split("\n") if line.strip()])
        
        st.success(f"✅ BPMN diagram generated with {len(elements)} elements!")
        
        return pretty_xml
        
    except Exception as e:
        st.error(f"❌ Error generating BPMN: {str(e)}")
        raise e

def clone_voice_elevenlabs(voice_file, voice_name: str, api_key: str) -> str:
    """Clone a voice using ElevenLabs API and return the voice ID."""
    try:
        client = ElevenLabs(api_key=api_key)
        
        st.info("🔄 Processing and uploading voice sample to ElevenLabs...")
        
        # Get the original file extension
        file_extension = os.path.splitext(voice_file.name)[1].lower()
        
        # Save uploaded file temporarily
        temp_dir = tempfile.gettempdir()
        temp_input_path = os.path.join(temp_dir, f"voice_input_{uuid.uuid4()}{file_extension}")
        temp_mp3_path = os.path.join(temp_dir, f"voice_sample_{uuid.uuid4()}.mp3")
        
        # Reset file pointer and save uploaded file
        voice_file.seek(0)
        file_bytes = voice_file.read()
        
        with open(temp_input_path, "wb") as f:
            f.write(file_bytes)
        
        # Verify input file was written correctly
        if not os.path.exists(temp_input_path) or os.path.getsize(temp_input_path) == 0:
            raise Exception("Failed to save uploaded audio file")
        
        st.info(f"📁 File uploaded: {voice_file.name} ({os.path.getsize(temp_input_path)} bytes)")
        
        # Convert to MP3 format for better ElevenLabs compatibility
        try:
            st.info("🔄 Converting to MP3 format (recommended by ElevenLabs)...")
            audio = AudioSegment.from_file(temp_input_path)
            
            # Check audio duration
            duration_seconds = len(audio) / 1000.0
            st.info(f"📊 Audio duration: {duration_seconds:.1f} seconds")
            
            # Warn if audio is too short
            if duration_seconds < 10:
                st.warning(f"⚠️ Audio is only {duration_seconds:.1f}s long. ElevenLabs recommends at least 30 seconds for best results.")
            
            # Export as MP3 with 192 kbps bitrate (ElevenLabs recommended)
            audio.export(
                temp_mp3_path,
                format="mp3",
                bitrate="192k",
                parameters=["-ar", "44100", "-ac", "1"]  # 44.1kHz, mono
            )
            
            st.success(f"✅ Converted to MP3: {os.path.getsize(temp_mp3_path)} bytes")
            
            # Use the converted MP3 file
            final_path = temp_mp3_path
            
        except Exception as e:
            st.warning(f"⚠️ Conversion failed ({str(e)}), using original file...")
            # If conversion fails, use original file
            final_path = temp_input_path
        
        # Clone the voice using IVC (Instant Voice Cloning)
        # Try with file handle first, then with file path if that fails
        st.info("🚀 Submitting to ElevenLabs API...")
        try:
            with open(final_path, "rb") as audio_file:
                voice = client.voices.ivc.create(
                    name=voice_name,
                    files=[audio_file],
                    description="Cloned voice for text-to-speech"
                )
        except Exception as file_handle_error:
            # If file handle fails, try with file path
            st.warning("⚠️ Retrying with alternative method...")
            voice = client.voices.ivc.create(
                name=voice_name,
                files=[final_path],
                description="Cloned voice for text-to-speech"
            )
        
        # Clean up temp files
        if os.path.exists(temp_input_path):
            os.remove(temp_input_path)
        if os.path.exists(temp_mp3_path):
            os.remove(temp_mp3_path)
        
        st.success(f"✅ Voice cloned successfully! Voice ID: {voice.voice_id}")
        return voice.voice_id
        
    except Exception as e:
        error_msg = str(e)
        
        # Check for common ElevenLabs errors and provide helpful messages
        if "can_not_use_instant_voice_cloning" in error_msg or "no access to use instant voice cloning" in error_msg:
            st.error("❌ **Subscription Required**: Your ElevenLabs plan doesn't include Instant Voice Cloning.")
            st.info("💡 **To use Voice Cloning**: Upgrade your ElevenLabs subscription at https://elevenlabs.io/pricing")
            st.info("📝 **Note**: Voice cloning requires a paid ElevenLabs plan (Starter or higher)")
        elif "invalid_content" in error_msg or "corrupted" in error_msg.lower():
            st.error("❌ **Invalid Audio File**: ElevenLabs rejected the audio file.")
            st.warning("🔍 **Possible Causes:**")
            st.info("   1. **Audio Quality Issues**: The audio may have excessive noise, distortion, or processing")
            st.info("   2. **Account Limitations**: Your ElevenLabs plan might have additional restrictions")
            st.info("   3. **Audio Characteristics**: Try using:")
            st.info("      • Clear, single-speaker audio")
            st.info("      • At least 30 seconds duration")
            st.info("      • Minimal background noise")
            st.info("      • No music or sound effects")
            st.markdown("---")
            st.info("💡 **Troubleshooting Steps:**")
            st.info("   1. Try recording new audio with your microphone in a quiet room")
            st.info("   2. Check your ElevenLabs account at https://elevenlabs.io/app/voice-lab")
            st.info("   3. Verify your subscription includes Instant Voice Cloning")
        elif "401" in error_msg or "Unauthorized" in error_msg:
            st.error("❌ **Invalid API Key**: Please check your ElevenLabs API key")
        elif "quota" in error_msg.lower():
            st.error("❌ **Quota Exceeded**: You've reached your ElevenLabs usage limit for this month")
        else:
            st.error(f"❌ **Error cloning voice**: {error_msg}")
        
        raise e

def generate_speech_elevenlabs(
    text: str,
    voice_id: str,
    api_key: str,
    stability: float = 0.5,
    similarity_boost: float = 0.75,
    style: float = 0.0,
    use_speaker_boost: bool = True
) -> str:
    """Generate speech using ElevenLabs with cloned voice."""
    try:
        client = ElevenLabs(api_key=api_key)
        
        st.info("🔄 Generating speech with your cloned voice...")
        
        # Generate audio with voice settings
        audio = client.text_to_speech.convert(
            voice_id=voice_id,
            text=text,
            model_id="eleven_multilingual_v2",
            voice_settings=VoiceSettings(
                stability=stability,
                similarity_boost=similarity_boost,
                style=style,
                use_speaker_boost=use_speaker_boost
            )
        )
        
        # Save audio to temporary file
        temp_dir = tempfile.gettempdir()
        audio_path = os.path.join(temp_dir, f"voice_clone_{uuid.uuid4()}.mp3")
        
        with open(audio_path, "wb") as f:
            for chunk in audio:
                f.write(chunk)
        
        st.success("✅ Speech generated successfully!")
        return audio_path
        
    except Exception as e:
        st.error(f"❌ Error generating speech: {str(e)}")
        raise e

def encode_image_to_base64(image_file) -> str:
    """Encode uploaded or captured image to base64 string."""
    try:
        if isinstance(image_file, str):
            # If it's already a file path
            with open(image_file, "rb") as f:
                return base64.b64encode(f.read()).decode('utf-8')
        else:
            # If it's an UploadedFile object (from st.file_uploader or st.camera_input)
            # Both return UploadedFile objects in Streamlit
            return base64.b64encode(image_file.getvalue()).decode('utf-8')
    except Exception as e:
        raise Exception(f"Failed to encode image: {str(e)}")

async def image_chat(
    query: str,
    image_base64: str,
    model: str,
    openai_api_key: str,
    voice: str,
    chat_history: List[Dict]
) -> Dict:
    """Process image chat query with OpenAI Vision API."""
    try:
        st.info(f"🔄 Step 1: Analyzing image with {model}...")
        
        # Only GPT-4o and GPT-4-turbo support vision
        vision_models = ["gpt-4o", "gpt-4o-mini", "gpt-4-turbo"]
        if model not in vision_models and "via Groq" not in model:
            st.warning(f"⚠️ {model} doesn't support vision. Switching to gpt-4o")
            actual_model = "gpt-4o"
        elif "via Groq" in model:
            st.warning("⚠️ Open source models via Groq not supported for vision. Using gpt-4o")
            actual_model = "gpt-4o"
        else:
            actual_model = model
        
        # Setup OpenAI client
        client = OpenAI(api_key=openai_api_key)
        
        # Build messages with image
        messages = []
        
        # Add chat history if exists
        if chat_history:
            for msg in chat_history[-3:]:  # Last 3 exchanges
                messages.append({"role": msg["role"], "content": msg["content"]})
        
        # Add current query with image
        content = [
            {
                "type": "text",
                "text": query
            },
            {
                "type": "image_url",
                "image_url": {
                    "url": f"data:image/jpeg;base64,{image_base64}"
                }
            }
        ]
        
        messages.append({
            "role": "user",
            "content": content
        })
        
        st.info("🔄 Step 2: Generating response...")
        response = client.chat.completions.create(
            model=actual_model,
            messages=messages,
            max_tokens=1000
        )
        
        text_response = response.choices[0].message.content
        if not text_response:
            raise Exception("No response generated from the model")
        
        st.write(f"Generated response of length: {len(text_response)}")
        
        # Generate voice
        st.info("🔄 Step 3: Generating audio...")
        async_openai = AsyncOpenAI(api_key=openai_api_key)
        
        try:
            audio_response = await async_openai.audio.speech.create(
                model="tts-1",
                voice=voice,
                input=text_response
            )
            
            temp_dir = tempfile.gettempdir()
            audio_path = os.path.join(temp_dir, f"response_{uuid.uuid4()}.mp3")
            
            with open(audio_path, "wb") as f:
                f.write(audio_response.content)
            st.success(f"✅ MP3 file generated successfully")
            
        except Exception as audio_error:
            st.warning(f"⚠️ Audio generation failed: {str(audio_error)}")
            audio_path = None
        
        st.success("✅ Image analysis complete!")
        return {
            "status": "success",
            "text_response": text_response,
            "audio_path": audio_path,
            "model_used": actual_model
        }
    
    except Exception as e:
        st.error(f"❌ Error during image chat processing: {str(e)}")
        return {
            "status": "error",
            "error": str(e),
            "query": query
        }

async def general_chat(
    query: str,
    model: str,
    openai_api_key: str,
    voice: str,
    chat_history: List[Dict]
) -> Dict:
    """Process general chat query without document context."""
    try:
        st.info(f"🔄 Step 1: Processing with {model}...")
        
        # Determine if using OpenAI or open source model
        if "via Groq" in model:
            st.warning("⚠️ Open source models via Groq require Groq API key (not implemented yet). Falling back to OpenAI GPT-4o-mini")
            actual_model = "gpt-4o-mini"
        else:
            actual_model = model
        
        # Build messages with chat history for direct API call
        messages = []
        if chat_history:
            # Add last 3 exchanges from history
            for msg in chat_history[-3:]:
                messages.append({"role": msg["role"], "content": msg["content"]})
        
        # Add current user query
        messages.append({"role": "user", "content": query})
        
        st.info("🔄 Step 2: Generating response...")
        # Use OpenAI API directly instead of Agent/Runner
        client = OpenAI(api_key=openai_api_key)
        response = client.chat.completions.create(
            model=actual_model,
            messages=messages,
            max_tokens=1000
        )
        
        text_response = response.choices[0].message.content
        if not text_response:
            raise Exception("No response generated from the model")
        
        st.write(f"Generated response of length: {len(text_response)}")
        
        # Generate voice
        st.info("🔄 Step 3: Generating audio...")
        async_openai = AsyncOpenAI(api_key=openai_api_key)
        
        try:
            # Try real-time playback
            try:
                from openai.helpers import LocalAudioPlayer
                async with async_openai.audio.speech.with_streaming_response.create(
                    model="gpt-4o-mini-tts",
                    voice=voice,
                    input=text_response,
                    response_format="pcm",
                ) as stream_response:
                    st.write("🎵 Playing audio (real-time)...")
                    await LocalAudioPlayer().play(stream_response)
                    st.success("✅ Audio playback complete")
            except (ImportError, OSError) as e:
                st.warning(f"⚠️ Real-time audio playback not available")
                st.info("💡 Generating downloadable MP3 instead...")
            
            # Generate MP3
            st.write("Generating downloadable MP3...")
            audio_response = await async_openai.audio.speech.create(
                model="gpt-4o-mini-tts",
                voice=voice,
                input=text_response,
                response_format="mp3"
            )
            
            temp_dir = tempfile.gettempdir()
            audio_path = os.path.join(temp_dir, f"response_{uuid.uuid4()}.mp3")
            
            with open(audio_path, "wb") as f:
                f.write(audio_response.content)
            st.success(f"✅ MP3 file generated successfully")
            
        except Exception as audio_error:
            st.warning(f"⚠️ Audio generation failed: {str(audio_error)}")
            audio_path = None
        
        st.success("✅ Chat processing complete!")
        return {
            "status": "success",
            "text_response": text_response,
            "audio_path": audio_path,
            "model_used": actual_model
        }
    
    except Exception as e:
        st.error(f"❌ Error during chat processing: {str(e)}")
        return {
            "status": "error",
            "error": str(e),
            "query": query
        }

def main() -> None:
    """Main application function."""
    st.set_page_config(
        page_title="AI Voice Assistant",
        page_icon="🤖",
        layout="wide",
        initial_sidebar_state="expanded"
    )
    
    # Custom CSS for better UI
    st.markdown("""
        <style>
        /* Main container styling */
        .main {
            padding: 2rem;
        }
        
        /* Title styling */
        h1 {
            color: #1f77b4;
            font-weight: 600;
            margin-bottom: 1rem;
        }
        
        /* Info box styling */
        .stAlert {
            border-radius: 10px;
            padding: 1rem;
        }
        
        /* Input fields */
        .stTextInput > div > div > input {
            border-radius: 8px;
        }
        
        /* Buttons */
        .stButton > button {
            border-radius: 8px;
            padding: 0.5rem 2rem;
            font-weight: 500;
        }
        
        /* File uploader */
        .stFileUploader {
            border: 2px dashed #4CAF50;
            border-radius: 10px;
            padding: 1rem;
        }
        
        /* Sidebar styling */
        .css-1d391kg {
            padding: 2rem 1rem;
        }
        
        /* Response box */
        .response-box {
            background-color: #f8f9fa;
            border-radius: 10px;
            padding: 1.5rem;
            margin: 1rem 0;
        }
        
        /* Audio player */
        .stAudio {
            margin: 1rem 0;
        }
        </style>
    """, unsafe_allow_html=True)
    
    init_session_state()
    setup_sidebar()
    
    # Dynamic title and description based on mode
    if st.session_state.mode == "General Chat":
        st.title("🤖 AI Chat Assistant")
        st.info("Ask me anything! Get AI-powered answers with voice responses using your selected model.")
    elif st.session_state.mode == "Image Chat":
        st.title("📸 AI Vision Chat")
        st.info("Upload or capture an image, then ask questions about it! Get AI-powered visual analysis with voice responses.")
    elif st.session_state.mode == "Process Flow Creator":
        st.title("🔄 Process Flow Creator")
        st.info("Generate professional swimlane diagrams from process descriptions! AI-powered process mapping in PowerPoint and Excel formats.")
    elif st.session_state.mode == "Voice Cloning":
        st.title("🎤 Voice Cloning & Text-to-Speech")
        st.info("Clone any voice and convert text to speech with advanced quality controls! Upload a voice sample and generate realistic speech.")
    else:
        st.title("🎙️ Voice RAG Agent")
        st.info("Get voice-powered answers to your documentation questions by uploading PDF documents and asking questions!")
    
    # File upload section - only for RAG mode
    if st.session_state.mode == "RAG (Document Q&A)":
        uploaded_file = st.file_uploader("Upload PDF", type=["pdf"])
        
        if uploaded_file:
            file_name = uploaded_file.name
            if file_name not in st.session_state.processed_documents:
                with st.spinner('Processing PDF...'):
                    try:
                        # Setup Qdrant if not already done
                        if not st.session_state.client:
                            client, embedding_model = setup_qdrant()
                            st.session_state.client = client
                            st.session_state.embedding_model = embedding_model
                        
                        # Process and store document
                        documents = process_pdf(uploaded_file)
                        if documents:
                            store_embeddings(
                                st.session_state.client,
                                st.session_state.embedding_model,
                                documents,
                                COLLECTION_NAME
                            )
                            st.session_state.processed_documents.append(file_name)
                            st.success(f"✅ Added PDF: {file_name}")
                            st.session_state.setup_complete = True
                    except Exception as e:
                        st.error(f"Error processing document: {str(e)}")
        
        # Display processed documents
        if st.session_state.processed_documents:
            st.sidebar.header("📚 Processed Documents")
            for doc in st.session_state.processed_documents:
                st.sidebar.text(f"📄 {doc}")
    
    # Image upload/capture section - only for Image Chat mode
    if st.session_state.mode == "Image Chat":
        st.markdown("### 📸 Upload or Capture Image")
        
        # Create two columns for upload and camera
        col1, col2 = st.columns(2)
        
        with col1:
            st.markdown("#### 📁 Upload Image")
            uploaded_image = st.file_uploader(
                "Choose an image file",
                type=["png", "jpg", "jpeg", "webp"],
                help="Upload an image to analyze"
            )
            
            if uploaded_image:
                st.session_state.current_image = uploaded_image
                st.session_state.image_base64 = encode_image_to_base64(uploaded_image)
                st.image(uploaded_image, caption="Uploaded Image", use_container_width=True)
        
        with col2:
            st.markdown("#### 📷 Capture from Camera")
            camera_image = st.camera_input("Take a photo")
            
            if camera_image:
                st.session_state.current_image = camera_image
                st.session_state.image_base64 = encode_image_to_base64(camera_image)
        
        # Show image info if available
        if st.session_state.current_image:
            st.success("✅ Image loaded! You can now ask questions about it.")
    
    # Process Flow Creator section
    if st.session_state.mode == "Process Flow Creator":
        st.info("💡 **Tip**: Describe your process with roles/departments involved. AI will create a professional SWIMLANE diagram showing who does what!")
        
        st.markdown("### 📝 Process Description")
        st.markdown("Describe your process with roles, activities, tools, and handoffs:")
        
        process_description = st.text_area(
            "Process details",
            placeholder="""Example:
Process: Employee Leave Request and Approval

Roles involved: Employee, Manager, HR Department

Process flow:
1. Employee submits leave request form (Tool: HR Portal)
2. Manager receives notification and reviews request
3. Manager decides: Approve or Reject?
   - If approved: Send to HR Department
   - If rejected: Notify employee
4. HR Department verifies leave balance (Tool: HRIS System)
5. HR Department updates leave calendar (Tool: Calendar System)
6. HR Department sends confirmation email to employee
7. End

Include tools, systems, and decision points!""",
            height=300,
            help="Describe the process steps, roles/departments, tools, and decision points. AI will organize them into swimlane diagram."
        )
        
        st.markdown("---")
        st.markdown("### 📌 What to Include for Best Results:")
        st.markdown("""
        - **Roles/Departments**: Who is responsible? (e.g., "Employee", "Manager", "HR", "System")
        - **Activities**: What tasks do they perform? (e.g., "Review request", "Approve")
        - **Tools/Systems**: What technology is used? (e.g., "CRM", "ERP System", "Database")
        - **Decision Points**: What choices are made? (e.g., "Is approved?", "In stock?")
        - **Handoffs**: When does work move between roles?
        """)
        
        st.markdown("**💡 Pro tip**: The AI will automatically detect roles and create horizontal swimlanes for each, showing cross-functional handoffs!")
        
        # Generate button
        if st.button("🎨 Generate Process Flow", disabled=not (st.session_state.openai_api_key and process_description)):
            with st.spinner("Generating process flow diagrams..."):
                try:
                    # Generate PowerPoint
                    pptx_path = asyncio.run(generate_process_flow(
                        process_description,
                        st.session_state.openai_api_key
                    ))
                    st.session_state.generated_flow_path = pptx_path
                    
                    # Extract flow_data for Excel generation
                    # Re-run AI analysis to get the flow data
                    client = OpenAI(api_key=st.session_state.openai_api_key)
                    prompt = f"""Analyze the following process description and create a structured SWIMLANE process flow diagram.

PROCESS DESCRIPTION:
{process_description}

Extract:
1. Process title/name
2. All roles/departments involved (swimlanes)
3. All steps in sequence with role assignments
4. Flow connections between steps (including cross-lane handoffs)

Format your response as JSON with this structure:
{{
    "title": "Process Name",
    "lanes": [
        {{
            "id": "lane1",
            "name": "Role/Department Name"
        }}
    ],
    "steps": [
        {{
            "id": 1,
            "lane_id": "lane1",
            "type": "start|process|decision|end",
            "text": "Step description",
            "tool": "Tool/System name (if applicable)",
            "next": [2]
        }}
    ]
}}

Step types:
- "start": Beginning of process (oval)
- "process": Regular activity/task (rectangle)
- "decision": Decision point with Yes/No branches (diamond)
- "end": End of process (oval)

Important:
- Identify ALL roles/departments involved and create lanes for each
- Assign EVERY step to the appropriate lane_id
- Include tools/systems used in each step
- Show handoffs between roles/departments through connections
- For decision points, specify both branches in the "next" array

Example roles: Employee, Manager, HR Department, System, Customer, etc.
"""
                    
                    response = client.chat.completions.create(
                        model="gpt-4o",
                        messages=[{"role": "user", "content": prompt}],
                        response_format={"type": "json_object"}
                    )
                    
                    import json
                    flow_data = json.loads(response.choices[0].message.content or "{}")
                    
                    # Generate Excel
                    excel_path = asyncio.run(generate_process_flow_excel(flow_data))
                    st.session_state.generated_excel_path = excel_path
                    
                except Exception as e:
                    st.error(f"Failed to generate process flow: {str(e)}")
        
        # Download section
        if st.session_state.generated_flow_path and os.path.exists(st.session_state.generated_flow_path):
            st.markdown("---")
            st.markdown("### ⬇️ Download Your Process Flow Diagrams")
            
            col1, col2 = st.columns(2)
            
            with col1:
                with open(st.session_state.generated_flow_path, "rb") as f:
                    flow_bytes = f.read()
                    st.download_button(
                        label="📥 Download PowerPoint (.pptx)",
                        data=flow_bytes,
                        file_name="process_flow_diagram.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        use_container_width=True
                    )
            
            with col2:
                if st.session_state.generated_excel_path and os.path.exists(st.session_state.generated_excel_path):
                    with open(st.session_state.generated_excel_path, "rb") as f:
                        excel_bytes = f.read()
                        st.download_button(
                            label="📥 Download Excel (.xlsx)",
                            data=excel_bytes,
                            file_name="process_flow_diagram.xlsx",
                            mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                            use_container_width=True
                        )
            
            st.success("✅ Process flow diagrams ready! Download in PowerPoint or Excel format (no gridlines, neat layout).")
    
    # Voice Cloning section
    if st.session_state.mode == "Voice Cloning":
        st.info("🎤 **Voice Cloning**: Upload a voice sample, then convert any text to speech using that cloned voice!")
        st.warning("⚠️ **Important**: Voice cloning requires a paid ElevenLabs subscription (Starter plan or higher). Free tier doesn't include this feature.")
        
        st.markdown("### 🎙️ Step 1: Upload Voice Sample")
        st.markdown("Upload an audio file of the voice you want to clone (at least 30 seconds recommended)")
        
        voice_file = st.file_uploader(
            "Upload Voice Sample",
            type=["mp3", "wav", "m4a", "flac"],
            help="Upload a clear audio recording of the voice you want to clone"
        )
        
        if voice_file:
            st.audio(voice_file, format=f"audio/{voice_file.type.split('/')[-1]}")
            st.success("✅ Voice sample uploaded!")
            
            voice_name = st.text_input(
                "Voice Name",
                placeholder="e.g., My Voice, John's Voice",
                help="Give your cloned voice a name"
            )
            
            if st.button("🧬 Clone Voice", disabled=not (st.session_state.elevenlabs_api_key and voice_name)):
                try:
                    voice_id = clone_voice_elevenlabs(
                        voice_file,
                        voice_name,
                        st.session_state.elevenlabs_api_key
                    )
                    st.session_state.cloned_voice_id = voice_id
                    st.session_state.voice_file = voice_file
                except Exception as e:
                    st.error(f"Failed to clone voice: {str(e)}")
        
        st.markdown("---")
        st.markdown("### 📝 Step 2: Enter Text to Convert")
        
        if st.session_state.cloned_voice_id:
            st.success(f"✅ Voice cloned! ID: {st.session_state.cloned_voice_id}")
            
            text_to_convert = st.text_area(
                "Text to Convert to Speech",
                placeholder="Enter the text you want to convert to speech using your cloned voice...",
                height=200,
                help="Type or paste the text you want to hear in your cloned voice"
            )
            
            if st.button("🔊 Generate Speech", disabled=not text_to_convert):
                try:
                    audio_path = generate_speech_elevenlabs(
                        text=text_to_convert,
                        voice_id=st.session_state.cloned_voice_id,
                        api_key=st.session_state.elevenlabs_api_key,
                        stability=st.session_state.voice_stability,
                        similarity_boost=st.session_state.voice_similarity,
                        style=st.session_state.voice_style,
                        use_speaker_boost=st.session_state.use_speaker_boost
                    )
                    
                    st.markdown("---")
                    st.markdown("### 🎧 Generated Audio")
                    
                    # Play audio
                    with open(audio_path, "rb") as audio_file:
                        audio_bytes = audio_file.read()
                        st.audio(audio_bytes, format="audio/mp3")
                    
                    # Download button
                    st.download_button(
                        label="📥 Download Audio (MP3)",
                        data=audio_bytes,
                        file_name="cloned_voice_output.mp3",
                        mime="audio/mp3"
                    )
                    
                except Exception as e:
                    st.error(f"Failed to generate speech: {str(e)}")
        else:
            st.info("👆 Please clone a voice first before generating speech")
        
        st.markdown("---")
        st.markdown("### 💡 Tips for Best Results")
        st.markdown("""
        - **Voice Sample Quality**: Use clear audio with minimal background noise
        - **Sample Length**: At least 30 seconds for better cloning (1-2 minutes ideal)
        - **Stability**: Higher values = more consistent, lower = more expressive
        - **Similarity Boost**: How closely to match the original voice
        - **Style**: Amplify the speaking style and emotion in the voice
        - **Speaker Boost**: Enhances clarity - keep enabled for best quality
        """)
    
    # BPMN Diagram Generator section
    if st.session_state.mode == "BPMN Diagram Generator":
        st.info("📊 **BPMN Diagram Generator**: Describe your business process and generate a standard BPMN XML diagram that can be imported into any BPMN tool!")
        
        st.markdown("### 📝 Process Description")
        st.markdown("Describe your business process with activities, decisions, and flow:")
        
        bpmn_description = st.text_area(
            "Business Process Details",
            placeholder="""Example:
Process: Customer Order Processing

1. Customer places order online (start)
2. System validates order details
3. Check inventory: Is product in stock?
   - Yes: Process payment
   - No: Notify customer (end)
4. Process payment: Is payment successful?
   - Yes: Create shipping label
   - No: Cancel order (end)
5. Warehouse picks and packs items
6. Ship order to customer
7. Send tracking information email
8. Order complete (end)

Include:
- Start and end points
- Tasks and activities
- Decision points (gateways)
- System vs human tasks
- Parallel activities (if any)""",
            height=350,
            help="Describe the business process with clear steps, decisions, and flow logic. AI will generate standard BPMN XML."
        )
        
        st.markdown("---")
        st.markdown("### 📌 What is BPMN?")
        st.markdown("""
        **BPMN (Business Process Model and Notation)** is an international standard for business process diagrams.
        
        **Supported Elements:**
        - **Start Event**: Beginning of the process ⚪
        - **End Event**: Completion of the process ⚫
        - **Tasks**: Activities performed (User Task, Service Task, etc.) ▭
        - **Gateways**: Decision points and parallel flows ◇
          - Exclusive: Choose one path (XOR)
          - Parallel: All paths simultaneously (AND)
          - Inclusive: One or more paths (OR)
        
        **Generated BPMN XML can be imported into:**
        - Camunda Modeler
        - bpmn.io
        - Signavio
        - Bizagi
        - Any BPMN 2.0 compatible tool
        """)
        
        # Generate button
        if st.button("🎨 Generate BPMN Diagram", disabled=not (st.session_state.openai_api_key and bpmn_description)):
            with st.spinner("Generating BPMN diagram..."):
                try:
                    bpmn_xml = asyncio.run(generate_bpmn_diagram(
                        bpmn_description,
                        st.session_state.openai_api_key
                    ))
                    st.session_state.generated_bpmn_xml = bpmn_xml
                    
                except Exception as e:
                    st.error(f"Failed to generate BPMN diagram: {str(e)}")
        
        # Download section
        if st.session_state.generated_bpmn_xml:
            st.markdown("---")
            st.markdown("### ⬇️ Download Your BPMN Diagram")
            
            # Show XML preview
            with st.expander("👀 Preview BPMN XML"):
                st.code(st.session_state.generated_bpmn_xml, language="xml")
            
            # Download button
            st.download_button(
                label="📥 Download BPMN XML (.bpmn)",
                data=st.session_state.generated_bpmn_xml,
                file_name="business_process_diagram.bpmn",
                mime="application/xml",
                use_container_width=True
            )
            
            st.success("✅ BPMN diagram ready! Download the .bpmn file and import it into your favorite BPMN tool.")
            
            st.markdown("---")
            st.markdown("### 🔧 How to Use the BPMN File")
            st.markdown("""
            1. **Download** the .bpmn file above
            2. **Open** your BPMN modeling tool:
               - [Camunda Modeler](https://camunda.com/download/modeler/) (Free, Desktop)
               - [bpmn.io](https://demo.bpmn.io/) (Free, Online)
               - Signavio, Bizagi, or any BPMN 2.0 tool
            3. **Import/Open** the downloaded file
            4. **Use Auto-Layout** (recommended for complex processes):
               - In Camunda Modeler: Right-click canvas → "Align Elements"
               - In bpmn.io: The diagram includes basic layout coordinates
            5. **Edit** the diagram visually with drag-and-drop
            6. **Export** to PNG, PDF, or keep as XML for execution
            
            💡 **Pro Tip**: The generated diagram includes basic horizontal layout. For complex processes with many branches, use your BPMN tool's auto-layout feature for optimal positioning!
            
            ⚠️ **Note**: Complex processes with multiple gateways and branches may need manual adjustment or auto-layout in your BPMN tool for best visual results.
            """)
    
    # Determine if ready to chat
    if st.session_state.mode == "General Chat":
        is_ready = bool(st.session_state.openai_api_key)
        placeholder = "e.g., What is machine learning?"
        input_label = "Ask me anything:"
    elif st.session_state.mode == "Image Chat":
        is_ready = bool(st.session_state.openai_api_key and st.session_state.image_base64)
        placeholder = "e.g., What do you see in this image?"
        input_label = "Ask about the image:"
    elif st.session_state.mode == "Process Flow Creator":
        # Process Flow mode doesn't use the chat interface
        is_ready = False
        placeholder = ""
        input_label = ""
    elif st.session_state.mode == "Voice Cloning":
        # Voice Cloning mode doesn't use the chat interface
        is_ready = False
        placeholder = ""
        input_label = ""
    elif st.session_state.mode == "BPMN Diagram Generator":
        # BPMN mode doesn't use the chat interface
        is_ready = False
        placeholder = ""
        input_label = ""
    else:
        is_ready = st.session_state.setup_complete
        placeholder = "e.g., How do I authenticate API requests?"
        input_label = "What would you like to know about the documentation?"
    
    # Query interface
    query = st.text_input(
        input_label,
        placeholder=placeholder,
        disabled=not is_ready
    )
    
    if query and is_ready:
        with st.status("Processing your query...", expanded=True) as status:
            try:
                # Route to appropriate function based on mode
                if st.session_state.mode == "General Chat":
                    result = asyncio.run(general_chat(
                        query,
                        st.session_state.selected_model,
                        st.session_state.openai_api_key,
                        st.session_state.selected_voice,
                        st.session_state.chat_history
                    ))
                    
                    # Update chat history
                    st.session_state.chat_history.append({"role": "user", "content": query})
                    if result["status"] == "success":
                        st.session_state.chat_history.append({"role": "assistant", "content": result["text_response"]})
                
                elif st.session_state.mode == "Image Chat":
                    result = asyncio.run(image_chat(
                        query,
                        st.session_state.image_base64,
                        st.session_state.selected_model,
                        st.session_state.openai_api_key,
                        st.session_state.selected_voice,
                        st.session_state.image_chat_history
                    ))
                    
                    # Update image chat history
                    st.session_state.image_chat_history.append({"role": "user", "content": query})
                    if result["status"] == "success":
                        st.session_state.image_chat_history.append({"role": "assistant", "content": result["text_response"]})
                
                else:
                    result = asyncio.run(process_query(
                        query,
                        st.session_state.client,
                        st.session_state.embedding_model,
                        COLLECTION_NAME,
                        st.session_state.openai_api_key,
                        st.session_state.selected_voice
                    ))
                
                if result["status"] == "success":
                    status.update(label="✅ Query processed!", state="complete")
                    
                    st.markdown("### Response:")
                    st.write(result["text_response"])
                    
                    # Show model used in General Chat and Image Chat modes
                    if (st.session_state.mode in ["General Chat", "Image Chat"]) and "model_used" in result:
                        st.caption(f"Model: {result['model_used']}")
                    
                    if "audio_path" in result and result["audio_path"]:
                        st.markdown(f"### 🔊 Audio Response (Voice: {st.session_state.selected_voice})")
                        try:
                            with open(result["audio_path"], "rb") as audio_file:
                                audio_bytes = audio_file.read()
                                st.audio(audio_bytes, format="audio/mp3", start_time=0)
                                st.download_button(
                                    label="📥 Download Audio Response",
                                    data=audio_bytes,
                                    file_name=f"voice_response_{st.session_state.selected_voice}.mp3",
                                    mime="audio/mp3"
                                )
                        except Exception as e:
                            st.error(f"❌ Error loading audio file: {str(e)}")
                    else:
                        st.info("💡 Audio generation was skipped due to environment limitations.")
                    
                    # Show sources only in RAG mode
                    if st.session_state.mode == "RAG (Document Q&A)" and "sources" in result:
                        st.markdown("### Sources:")
                        for source in result["sources"]:
                            st.markdown(f"- {source}")
                else:
                    status.update(label="❌ Error processing query", state="error")
                    st.error(f"Error: {result.get('error', 'Unknown error occurred')}")
            
            except Exception as e:
                status.update(label="❌ Error processing query", state="error")
                st.error(f"Error processing query: {str(e)}")
    
    elif not is_ready:
        if st.session_state.mode == "General Chat":
            st.info("👈 Please configure your OpenAI API key in the sidebar to start chatting!")
        else:
            st.info("👈 Please configure the system and upload documents first!")

if __name__ == "__main__":
    main()